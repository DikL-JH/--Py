print('                                                           《 <<<  欢迎使用运营数据处理系统  >>> 》            \n')

print('                                                                    系统正在加载中--请稍后\n\n')


#-*- coding: UTF-8 -*-
#ppt排序C:\\Users\\admin0001\\Downloads
import os
import shutil
import pptx
import random
import comtypes.client
from os import listdir
from time import sleep
from pptx import Presentation
from PIL import Image
path='C:\\ppt'  #ppt文件所在位置
path2 = 'c:\\kw.txt' #关键词文件所在位置
kg=0
width_i = 280
height_i = 157
row_max = 5
line_max = 2
all_path = list()
num = 0

#ppt批量转换为pptx
powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
files = listdir(path)
for file in files:
    try:
        if file.endswith(".ppt"):
            pp=path+'\\'+file
            print(pp)
            pr=powerpoint.Presentations.Open(pp,WithWindow=False)
            pr.SaveAs(pp[:-3]+'pptx')
            pr.Close()
            os.remove(pp)
    except Exception as e:
                print(e)
                print("**************************************")
                print("*数据异常--请手动处理--按回车继续处理*")
                input("**************************************")


#排查敏感字符
def text_r():
    ee=0
    files = listdir(path)
    for file in files:
        if file.endswith(".pptx"):
            try:
                ee=ee+1
                pp=path+'\\'+file
                print("打开文件  "+str(ee)+"    "+str(file))
                TEXT_NEED_REPLACE = [('201X','20XX'),('鹿大仙素材','XXX'),('奈森设计','XXX'),('2016','20XX'),('2017','20XX'),('2018','20XX'),('2019','20XX'),('1PPT',''),('情缘素材',''),
                             ('PPT模板网-WWW',''),('亮亮图文旗舰店','XXX'),('https://liangliangtuwen.tmall.com',''),
                             ('亮亮图文旗舰店https://liangliangtuwen.tmall.com',''),('W ww.51pptmoban.com',''),
                             ('WWW.1PPT.COM',''),('www.1ppt.com',''),('模板网',''),('1PPT.COM',''),('第一',''),('旗舰店',''),('大白素材',''),
                             ('读爱',''),('读爱礼坊',''),('米鸽米','XXX'),('米鸽设计','XXX'),('千图网',''),('凤凰办公',''),('亮亮图文','')]#需要替换的内容
                TEXT_NEED_REPLACE2 = [('版权声明'),('包图网')]
                #读取kw
                with open(path2, 'r',encoding='UTF-8') as f:
                    s = [i[:-1].split('，') for i in f.readlines()]
                for i in range(len(s)):
                    TEXT_NEED_REPLACE.extend([(s[i][0], s[i][1])])

                FILE_OPEN = pp#需要打开的文件地址
                
                FILE_SAVE = pp#需要保存的文件地址
                

                def replace_text(text_frame):#该函数实现的是文本替换功能
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            for tt in TEXT_NEED_REPLACE:
                                if tt[0] in run.text:
                                    #print("替换")
                                    run.text = run.text.replace(tt[0], tt[1])
                            for tt2 in TEXT_NEED_REPLACE2:
                                if tt2[0] in run.text:
                                    global kg
                                    kg=1
                                    #print("包图")

                def process_ppt(filename_open, filename_save):
                    prs = Presentation(filename_open)
                    #rId = prs.slides._sldIdLst[-1].rId
                    #prs.part.drop_rel(rId)
                    #del prs.slides._sldIdLst[-1]
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:#判断Shape是否含有文本框
                                text_frame = shape.text_frame
                                replace_text(text_frame)#调用replace_text函数实现文本替换
                                global kg
                                if kg==1 :
                                    rId = prs.slides._sldIdLst[-1].rId
                                    prs.part.drop_rel(rId)
                                    del prs.slides._sldIdLst[-1]
                                    kg=0
                               
                                
                            if shape.has_table:#判断Shape是否含有表格
                                table = shape.table
                                for cell in table.iter_cells():#遍历表格的cell
                                    text_frame = cell.text_frame
                                    replace_text(text_frame)#调用replace_text函数实现文本替换
                    prs.save(filename_save)#保存

                process_ppt(FILE_OPEN, FILE_SAVE)
            except Exception as e:
                print(e)
                print("**************************************")
                print("*数据异常--请手动处理--按回车继续处理*")
                input("**************************************")
        else:
            print('跳过 '+file+' 文件 ')
            continue



#内容图生成
def nrt():
    import comtypes.client
    import time
    from time import sleep

    def init_powerpoint():
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        return powerpoint
        
    def ppt_to_picture(powerpoint, inputFileName, formatType = 32):
        deck = powerpoint.Presentations.Open(inputFileName,WithWindow=False)
        deck.SaveAs(inputFileName.rsplit('.')[0] + '.JPG', 17)
        deck.Close()

    def convert_files_in_folder(powerpoint, folder):
        files = os.listdir(folder)
        pptfiles = [f for f in files if f.endswith((".ppt", ".pptx"))]
        sss=0
        for pptfile in pptfiles:
            sss=sss+1
            try:
                fullpath = os.path.join(path, pptfile)
                print('导出图片'+str(sss)+'   '+pptfile)
                ppt_to_picture(powerpoint, fullpath, fullpath)
            except Exception as e:
                print(e)
                print("**************************************")
                print("*数据异常--请手动处理--按回车继续处理*")
                input("**************************************")
            

    if __name__ == "__main__":
        powerpoint = init_powerpoint()
        convert_files_in_folder(powerpoint, path)
        powerpoint.Quit()


#封面图生成
def fmt():
    global num
    global all_path
    global height_i
    Paths=listdir(path)
    for  Paths in (Paths):
        if Paths.endswith((".ppt", ".pptx")):
            continue
        else:
            try:
                files=listdir(path+'\\'+Paths)
                files.sort(key= lambda x:str(x[:-5]))
                for i in range(len(files)):
                    if ".JPG" in files[i]:
                        paths=path+'\\'+Paths+'\\'+files[i]
                        all_path.append(paths)
                while True:
                    if len(all_path) < 9:
                        all_path.append(all_path[random.randint(int(len(all_path) / 2), len(all_path) - 1)])
                    else:
                        break
                toImage = Image.new('RGBA',(560,866),(228,228,228))
                toImage=toImage.convert('RGB')
                #print(all_path)
                for i in range(row_max):
                    for j in range(line_max):
                        pic_fole_head = Image.open(all_path[num])
                        w,h = pic_fole_head.size 
                        if i==0:
                            _h = int(h/w*560)
                            height_i = int(h/w*280)
                            _h2 = int(_h/2)+1
                            #print(int(_h))
                            #print(int(_h/2)+1)
                            tmppic = pic_fole_head.resize((560, _h))
                            toImage.paste(tmppic, (0,0))
                            loc=(0,0)
                            num=num+1
                            break
                        tmppic = pic_fole_head.resize((width_i, height_i))
                        if j==1:
                            loc = (int(j*width_i)+1,int(i*height_i)+_h2+i)
                        else:
                            loc = (int(j*width_i),int(i*height_i)+_h2+i)
                        
                        toImage.paste(tmppic, loc)
                        num=num+1
                print('%s\\%s.jpg'%(path+'\\'+Paths,Paths))        
                toImage.save('%s\\%s.jpg'%(path+'\\'+Paths,Paths))
            
            except Exception as e:
                print(e)
                print("**************************************")
                print("*数据异常--请手动处理--按回车继续处理*")
                input("**************************************")    
        all_path = list()
        num = 0


#视频生成
def mpf():
    # -*- coding: UTF-8 -*-
    import win32com.client
    import time
    import os
    import shutil
    import re
    from os import listdir
    import shutil

    def ppt_to_mp4(ppt_path,mp4_target,resolution = 480,frames = 20,quality = 30,timeout = 600):
        # status:Convert result. 0:failed. -1: timeout. 1:success.
        status = 0
        if ppt_path == '' or mp4_target == '':
            return status
        # start_tm:Start time
        start_tm = time.time()

        # Create a folder that does not exist.
        sdir = mp4_target[:mp4_target.rfind('\\')]
        if not os.path.exists(sdir):
            os.makedirs(sdir)

        # Start converting
        ppt = win32com.client.Dispatch('PowerPoint.Application')
        presentation = ppt.Presentations.Open(ppt_path,WithWindow=False)
        # CreateVideo() function usage: https://docs.microsoft.com/en-us/office/vba/api/powerpoint.presentation.createvideo
        presentation.CreateVideo(mp4_target,-1,1,resolution,frames,quality)
        while True:
            try:
                time.sleep(0.2)
                if time.time() - start_tm > timeout:
                    # Converting time out. Killing the PowerPoint process(An exception will be threw out).
                    os.system("taskkill /f /im POWERPNT.EXE")
                    status = -1
                    break
                if os.path.exists(mp4_path) and os.path.getsize(mp4_target) == 0:
                    # The filesize is 0 bytes when convert do not complete.
                    continue
                status = 1
                break
            except Exception as e:
                print("Error! Code: {c}, Message, {m}").format(c = type(e).__name__, m = str(e))
                
        print (time.time()-start_tm)
        if status != -1:
            ppt.Quit()

        return status
        
    if __name__ == '__main__':

        quality = 30
        resolution = 480
        #print(1)
        frames = 20
        ie_temp_dir = ''
        status = 0
        timeout = 10*60
        files = listdir(path)
        bbb=0
        for file in files:
            #print(2)
            #print(file)
            try:
                if file.endswith((".ppt", ".pptx")):
                    bbb=bbb+1
                    #print(path)
                    Path=path+'\\'+file
                    #print(Path)
                    #print(3)
                    ppt_path = os.path.abspath(Path)
                    pattern = re.compile(r'([^<>/\\\|:""\*\?]+)\.\w+$')
                    data = pattern.findall(Path)
                    print(str(data)+"   "+str(bbb))
                    mp4_path = os.path.abspath(path+'\\'+str(data[0])+'\\'+'1.mp4')
                    time.sleep(0.1)
                    try:
                        status = ppt_to_mp4(ppt_path,mp4_path,resolution,frames,quality,timeout)
                        if ie_temp_dir != '':
                            shutil.rmtree(ie_temp_dir, ignore_errors=True)
                    except Exception as e:
                        print ("Error! Code: {c}, Message, {m}".format(c = type(e).__name__, m = str(e)))
                                
                    if status == -1:
                        print ('Failed:timeout.')
                    elif status == 1:
                        print ('Success!')
                    else:
                        if os.path.exists(mp4_path):
                            os.remove(mp4_path)
                        print ('Failed:The ppt may have unknow elements. You can try to convert it manual.')
                    if '.pptx'in file:    
                        shutil.copyfile(Path,path+'\\'+str(data[0])+'\\'+str(data[0])+'.pptx')
                    elif '.ppt' in file:
                        shutil.copyfile(Path,path+'\\'+str(data[0])+'\\'+str(data[0])+'.ppt')
            except Exception as e:
                print(e)
                print("**************************************")
                print("*数据异常--请手动处理--按回车继续处理*")
                input("**************************************")                    

#ppt批量转换为pptx
def ppt():
    files = listdir(path)
    for file in files:
        try:
            if file.endswith(".ppt"):
                pp=path+'\\'+file
                print(pp)
                pr=powerpoint.Presentations.Open(pp,WithWindow=False)
                pr.SaveAs(pp[:-3]+'pptx')
                pr.Close()
                os.remove(pp)
        except Exception as e:
                print(e)
                print("**************************************")
                print("*数据异常--请手动处理--按回车继续处理*")
                input("**************************************")

def wj():
    files = listdir(path)
    for file in files:
        if file.endswith((".ppt", ".pptx")):
            continue
        else:
            Path = path+'\\'+file
            Files = listdir(Path)
            for File in Files:
                try:    
                    if File.endswith(".ppt"):
                        pp=Path+'\\'+File
                        print(pp)
                        pr=powerpoint.Presentations.Open(pp,WithWindow=False)
                        pr.SaveAs(pp[:-3]+'pptx')
                        pr.Close()
                        os.remove(pp)
                except Exception as e:
                        print(e)
                        print("**************************************")
                        print("*数据异常--请手动处理--按回车继续处理*")
                        input("**************************************")


kg='0'
while True:
    if kg==str(0):
        print('****************************************************************************************************************************************************')
        print('*   排查敏感字符    *    内容图生成     *     封面图生成     *    视频生成     *    一键生成所有文件    *     ppt批量转换为pptx      *      退出   *')
        print('*        1          *        2          *          3         *        4        *           5            *             6              *       T     *')
        kg=input('****************************************************************************************************************************************************\n')
    elif kg==str(1):
        text_r()
        kg='0'
        print('\n\n处理完成，请继续操作\n\n')
    elif kg==str(2):
        nrt()
        kg='0'
        print('\n\n处理完成，请继续操作\n\n')
    elif kg==str(3):
        kg1='0'
        while True:
            if kg1==str(0):
                print('***********************************************')
                print('*     ppt     *      文件夹      *     返回   *')
                print('*      1      *         2        *       3    *')
                kg1=input('***********************************************\n')
            elif kg1==str(1):
                nrt()
                fmt()
                kg1='0'
                print('\n\n处理完成，请继续操作\n\n')
            elif kg1==str(2):
                fmt()
                kg1='0'
                print('\n\n处理完成，请继续操作\n\n')
            elif kg1==str(3):
                kg='0'
                break
            else:
                kg1=input('输入有误，请重新输入\n')      
    elif kg==str(4):
        mpf()
        kg='0'
        print('\n\n处理完成，请继续操作\n\n')
    elif kg==str(5):
        nrt()
        fmt()
        mpf()
        kg='0'
        print('\n\n处理完成，请继续操作\n\n')
    elif kg == str(6):
        kg2 = '0'
        while True:
            if kg2 == str(0):
                print('***********************************')
                print('*   ppt   *    文件夹   *   退出  *')
                print('*    1    *      2       *    3   *')
                kg2 = input('***********************************\n')
            elif kg2 == str(1):
                ppt()
                kg2 = '0'
                print('\n\n处理完成，请继续操作\n\n')
            elif kg2 == str(2):
                wj()
                kg2 = '0'
                print('\n\n处理完成，请继续操作\n\n')
            elif kg2 == str(3):
                kg = '0'
                break
            else:
                kg2 = input('输入有误，请重新输入\n')
    elif kg==str('T'):
        break
    else:
        kg=input('输入有误，请重新输入\n')

